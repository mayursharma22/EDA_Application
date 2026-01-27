from typing import Optional
from typing import Dict
from typing import List
from typing import Tuple
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION


# ------------ utilities ------------
def _hex_to_rgb(hex_color: str) -> RGBColor:
    try:
        h = (hex_color or "").strip()
        if h.startswith("#") and len(h) == 7:
            r = int(h[1:3], 16)
            g = int(h[3:5], 16)
            b = int(h[5:7], 16)
            return RGBColor(r, g, b)
    except Exception:
        pass
    return RGBColor(0x12, 0x29, 0x5D)


def _norm(s: str) -> str:
    return (s or "").strip().lower()


def _ends_with_metric(col: str, metric: str) -> bool:
    return _norm(col).rsplit(" ", 1)[-1] == _norm(metric)


def _is_total_metric(col: str, metric: str) -> bool:
    return _norm(col) == _norm(f"total {metric}")


def _fmt_num(val, is_cost: bool) -> str:
    try:
        if pd.isna(val):
            return "-"
        x = float(val)
        return f"${x:,.0f}" if is_cost else f"{x:,.0f}"
    except Exception:
        return str(val)


def _to_list(series_like) -> List:
    if series_like is None:
        return []
    if isinstance(series_like, (pd.Series, pd.Index, np.ndarray)):
        return list(series_like)
    if isinstance(series_like, list):
        return series_like
    return list(series_like)


def _is_hex_color(s: str) -> bool:
    if not isinstance(s, str):
        return False
    s = s.strip()
    if not (s.startswith("#") and len(s) == 7):
        return False
    try:
        int(s[1:], 16)
        return True
    except ValueError:
        return False


def _add_table_in_rect(
    slide,
    df: pd.DataFrame,
    left,
    top,
    width,
    max_height,
    header_hex: str,
    min_body_row_h: int = Inches(0.22),
    header_row_h: int = Inches(0.28),
    prefer_rows: int | None = None,
) -> int:
    """
    Draw df as a PowerPoint table inside the given rectangle.
    Returns the actual height used (in EMUs).

    We keep the header row + as many body rows as can fit given min_body_row_h.
    """
    if df.empty or df.shape[1] == 0 or max_height <= 0:
        return 0

    cols = df.shape[1]

    remain_for_body = int(max_height - header_row_h)
    if remain_for_body <= 0:
        return 0

    max_body_fit = max(1, remain_for_body // int(min_body_row_h))
    if prefer_rows is not None:
        body_rows = min(prefer_rows, max_body_fit, len(df))
    else:
        body_rows = min(max_body_fit, len(df))

    rows = body_rows + 1
    used_height = header_row_h + int(body_rows * min_body_row_h)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, used_height)
    tbl = tbl_shape.table

    col_w = int(width / cols)
    for j in range(cols):
        tbl.columns[j].width = col_w

    header_rgb = _hex_to_rgb(header_hex)
    for j in range(cols):
        cell = tbl.cell(0, j)
        cell.text = str(df.columns[j])
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb
    tbl.rows[0].height = header_row_h

    for i in range(1, rows):
        for j in range(cols):
            try:
                val = df.iloc[i - 1, j]
            except Exception:
                val = "-"
            hdr = str(df.columns[j]).lower()
            is_cost = ("cost" in hdr) or ("spend" in hdr)
            cell = tbl.cell(i, j)
            cell.text = _fmt_num(val, is_cost)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
        tbl.rows[i].height = min_body_row_h

    return used_height


def _append_totals_like_excel(
    pivot: pd.DataFrame, metrics: List[str], weekly: bool
) -> pd.DataFrame:
    if pivot.empty or pivot.shape[1] <= 1:
        return pivot
    out = pivot.copy()
    first_col = out.columns[0]

    for m in metrics:
        m_cols = [
            c
            for c in out.columns[1:]
            if _ends_with_metric(c, m) and not _is_total_metric(c, m)
        ]
        if len(m_cols) > 1:
            tname = f"Total {m.capitalize()}"
            out[tname] = out[m_cols].sum(axis=1)

    non_total = [first_col] + [
        c for c in out.columns[1:] if not str(c).lower().startswith("total ")
    ]
    totals = [c for c in out.columns[1:] if str(c).lower().startswith("total ")]
    out = out[non_total + totals]

    for m in metrics:
        m_cols = [
            c
            for c in out.columns[1:]
            if _ends_with_metric(c, m) and not _is_total_metric(c, m)
        ]
        t_cols = [c for c in out.columns[1:] if _is_total_metric(c, m)]
        if len(m_cols) <= 1 and t_cols:
            out = out.drop(columns=t_cols, errors="ignore")

    if weekly:
        non_date = [c for c in out.columns if c != first_col]
        if len(non_date) == 1 and str(non_date[0]).lower().startswith("total "):
            out = out.drop(columns=[non_date[0]])

    return out


def _compute_pivots(df: pd.DataFrame, params: Dict) -> List[dict]:
    date_var = params.get("date_var", "date").strip()
    date_grain = (params.get("date_grain", "weekly") or "weekly").strip().lower()
    QC_variables: List[str] = [
        v.strip() for v in params.get("QC_variables", []) if v.strip()
    ]
    columns_breakdown: List[str] = [
        v.strip() for v in params.get("columns_breakdown", []) if v.strip()
    ]
    metrics_param: List[str] = [
        m.strip() for m in params.get("metrics", []) if m.strip()
    ]
    metric_var = params.get("metric_var", "metric").strip()
    value_var = params.get("value_var", "value").strip()
    week_start_day = (params.get("week_start_day", "") or "").strip()

    user_palette = params.get("graph_colors", []) or []
    graph_palette: List[str] = []
    _seen = set()
    for c in user_palette:
        cc = (c or "").strip().upper()
        if _is_hex_color(cc) and cc not in _seen:
            _seen.add(cc)
            graph_palette.append(cc)

    df = df.copy()
    df[date_var] = pd.to_datetime(df[date_var], errors="coerce")

    if date_grain == "weekly" and week_start_day:
        valid_days = [
            "Monday",
            "Tuesday",
            "Wednesday",
            "Thursday",
            "Friday",
            "Saturday",
            "Sunday",
        ]
        if week_start_day in valid_days:
            target_idx = valid_days.index(week_start_day)
            mask = df[date_var].notna()
            delta = (df.loc[mask, date_var].dt.weekday - target_idx) % 7
            df.loc[mask, date_var] = df.loc[mask, date_var] - pd.to_timedelta(
                delta, unit="D"
            )

    colmap = {c.strip().lower(): c for c in df.columns}
    has_long = (colmap.get(metric_var.lower()) in df.columns) and (
        colmap.get(value_var.lower()) in df.columns
    )

    if not has_long and metrics_param:
        wanted = {m.lower() for m in metrics_param}
        Metrics = [c for c in df.columns if c.strip().lower() in wanted]
    elif has_long:
        df[metric_var] = df[metric_var].astype(str).str.strip()
        Metrics = sorted(df[metric_var].dropna().unique().tolist())
    else:
        raise KeyError(
            "PPT generation requires long format (metric/value) OR explicit 'metrics' list."
        )

    import itertools

    if QC_variables:
        value_lists = []
        for v in QC_variables:
            vals = (
                pd.Series(df[v]).dropna().astype(str).map(lambda s: s.strip()).tolist()
            )
            vals = sorted(set(vals), key=str)
            value_lists.append(vals)
        QC_combos = list(itertools.product(*value_lists))
    else:
        QC_combos = [()]

    sections = []
    for combo in QC_combos:
        df_f = df.copy()
        for i, val in enumerate(combo):
            col = QC_variables[i]
            df_f = df_f[df_f[col].astype(str).str.strip() == str(val).strip()]
        if df_f.empty:
            continue
        valid_breakdown = [
            c
            for c in columns_breakdown
            if c in df_f.columns
            and df_f[c].dropna().astype(str).str.strip().ne("").any()
        ]

        group_cols = [date_var] + QC_variables + valid_breakdown

        if has_long:
            df_week_long = df_f.groupby(group_cols + [metric_var], as_index=False)[
                value_var
            ].sum()
            long = df_week_long.rename(
                columns={metric_var: "metric", value_var: "value"}
            )
        else:
            agg = {m: "sum" for m in Metrics}
            df_week = df_f.groupby(group_cols, as_index=False).agg(agg)
            id_vars = [date_var] + valid_breakdown
            long = df_week.melt(
                id_vars=id_vars,
                value_vars=list(Metrics),
                var_name="metric",
                value_name="value",
            )

        long["combined"] = (
            long[valid_breakdown]
            .fillna("")
            .astype(str)
            .agg(" ".join, axis=1)
            .str.replace(r"\s+", " ", regex=True)
            .str.strip()
            + " "
            + long["metric"].astype(str).str.strip()
            if valid_breakdown
            else long["metric"].astype(str).str.strip()
        )

        pivot_week = (
            long.groupby([date_var, "combined"], as_index=False)["value"]
            .sum()
            .pivot(index=date_var, columns="combined", values="value")
            .fillna(0)
            .reset_index()
        )
        pivot_week = _append_totals_like_excel(pivot_week, Metrics, weekly=True)

        df_f = df_f[df_f[date_var].notna()]
        df_f["Quarter"] = (
            df_f[date_var].dt.to_period("Q").astype(str).str.replace("Q", " Q")
        )
        q_group = ["Quarter"] + QC_variables + valid_breakdown

        if has_long:
            df_q_long = df_f.groupby(q_group + [metric_var], as_index=False)[
                value_var
            ].sum()
            q_long = df_q_long.rename(
                columns={metric_var: "metric", value_var: "value"}
            )
        else:
            df_q = df_f.groupby(q_group, as_index=False).agg(
                {m: "sum" for m in Metrics}
            )
            q_long = df_q.melt(
                id_vars=["Quarter"] + valid_breakdown,
                value_vars=list(Metrics),
                var_name="metric",
                value_name="value",
            )

        q_long["combined"] = (
            q_long[valid_breakdown]
            .fillna("")
            .astype(str)
            .agg(" ".join, axis=1)
            .str.replace(r"\s+", " ", regex=True)
            .str.strip()
            + " "
            + q_long["metric"].astype(str).str.strip()
            if valid_breakdown
            else q_long["metric"].astype(str).str.strip()
        )

        pivot_q = (
            q_long.groupby(["Quarter", "combined"], as_index=False)["value"]
            .sum()
            .pivot(index="Quarter", columns="combined", values="value")
            .fillna(0)
            .reset_index()
        )
        pivot_q = _append_totals_like_excel(pivot_q, Metrics, weekly=False)

        sheet_name = "_".join(map(str, combo)) if QC_variables else "General"

        sections.append(
            {
                "sheet_name": sheet_name[:31],
                "pivot_week": pivot_week,
                "pivot_q": pivot_q,
                "metrics": Metrics,
                "palette": graph_palette,
            }
        )

    return sections


def _build_q_summary_metric(pivot_q: pd.DataFrame, metric: str) -> pd.DataFrame:
    if pivot_q.empty or pivot_q.shape[1] == 0:
        return pivot_q.copy()

    cols = list(pivot_q.columns)
    quarter_col = cols[0]
    metric_cols = [c for c in cols[1:] if _ends_with_metric(c, metric)]
    if not metric_cols:
        return pivot_q[[quarter_col]].copy()

    total_cols = [c for c in metric_cols if _is_total_metric(c, metric)]
    non_total_cols = [c for c in metric_cols if not _is_total_metric(c, metric)]

    out = pd.DataFrame()
    out["Quarter"] = pivot_q[quarter_col]
    for c in non_total_cols:
        out[c] = pivot_q[c]
    for c in total_cols:
        out[c] = pivot_q[c]
    return out


def _align_cats_and_series(
    cats: List, series_list: List[List[float]]
) -> Tuple[List, List[List[float]]]:
    cats = _to_list(cats)
    lengths = [len(cats)] + [len(_to_list(s)) for s in series_list]
    if not lengths:
        return [], []
    m = min(lengths)
    cats = cats[:m]
    out = []
    for s in series_list:
        v = _to_list(s)[:m]
        v = (
            pd.to_numeric(pd.Series(v), errors="coerce")
            .fillna(0.0)
            .astype(float)
            .tolist()
        )
        out.append(v)
    return cats, out


def _add_metric_chart(
    slide,
    pivot_week: pd.DataFrame,
    metric: str,
    area_if_many=True,
    palette: List[str] | None = None,
    rect: Tuple[int, int, int, int] | None = None,
):
    """
    Add the metric chart inside the given rectangle (rect). If rect is None,
    it falls back to the old fixed coordinates. Chart title is disabled.
    """
    if pivot_week.empty:
        return None
    date_col = pivot_week.columns[0]
    dates = pd.to_datetime(pivot_week[date_col], errors="coerce")
    cats = dates.dt.strftime("%b '%y").fillna("").tolist()
    series_cols = [
        c
        for c in pivot_week.columns[1:]
        if _ends_with_metric(c, metric) and not _is_total_metric(c, metric)
    ]

    fallback_used = False
    if not series_cols:
        total_candidates = [
            c for c in pivot_week.columns if _is_total_metric(c, metric)
        ]
        if total_candidates:
            series_cols = [total_candidates[0]]
        else:
            metric_cols = [
                c for c in pivot_week.columns[1:] if _ends_with_metric(c, metric)
            ]
            if metric_cols:
                pivot_week = pivot_week.copy()
                pivot_week["__tmp_total__"] = pivot_week[metric_cols].sum(axis=1)
                series_cols = ["__tmp_total__"]
            else:
                return None
        fallback_used = True

    # --- Values + alignment ---
    series_values = [pivot_week[c].tolist() for c in series_cols]
    cats, series_values = _align_cats_and_series(cats, series_values)
    if not cats or not any(len(s) for s in series_values):
        return None

    # --- Chart type ---
    chart_type = (
        XL_CHART_TYPE.LINE
        if (not area_if_many or len(series_cols) <= 2)
        else XL_CHART_TYPE.AREA_STACKED
    )

    # --- Chart data ---
    chart_data = CategoryChartData()
    chart_data.categories = cats
    for name, vals in zip(series_cols, series_values):
        friendly = name
        if fallback_used and (
            name == "__tmp_total__" or _is_total_metric(name, metric)
        ):
            friendly = f"Total {metric.capitalize()}"
        chart_data.add_series(friendly, vals)

    # --- Position & size (use rect if provided) ---
    if rect:
        x, y, cx, cy = rect
    else:
        x, y, cx, cy = Inches(0.5), Inches(3.35), Inches(9.0), Inches(4.7)

    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

    # --- Styling ---
    chart.has_title = False
    try:
        chart.chart_style = 2
    except Exception:
        pass

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.has_major_gridlines = False

    pal = palette or []
    for idx, s in enumerate(getattr(chart, "series", [])):
        try:
            color_hex = pal[idx % len(pal)] if pal else None
            if color_hex and _is_hex_color(color_hex):
                rgb = _hex_to_rgb(color_hex)
                if chart_type == XL_CHART_TYPE.LINE:
                    s.format.line.width = Pt(2)
                    s.format.line.color.rgb = rgb
                else:
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = rgb
                    s.format.line.width = Pt(0.75)
                    s.format.line.color.rgb = rgb
        except Exception:
            pass

    return chart


def _add_metric_summary_table(slide, q3: pd.DataFrame, tab_header_hex: str):
    if q3.empty or q3.shape[1] == 0:
        return

    rows = min(len(q3) + 1, 12)
    cols = q3.shape[1]
    table_left, table_top = Inches(0.5), Inches(1.0)
    table_w, table_h = Inches(9.0), Inches(2.2)

    tbl_shape = slide.shapes.add_table(
        rows, cols, table_left, table_top, table_w, table_h
    )
    table = tbl_shape.table

    for j in range(cols):
        table.columns[j].width = int(table_w / cols)

    header_rgb = _hex_to_rgb(tab_header_hex)

    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = str(q3.columns[j])
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb

    for i in range(1, rows):
        for j in range(cols):
            try:
                val = q3.iloc[i - 1, j]
            except Exception:
                val = "-"
            hdr = str(q3.columns[j]).lower()
            is_cost = ("cost" in hdr) or ("spend" in hdr)
            cell = table.cell(i, j)
            cell.text = _fmt_num(val, is_cost)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT


def _add_metric_slide(
    pres,
    sheet_name: str,
    pivot_week: pd.DataFrame,
    pivot_q: pd.DataFrame,
    metric: str,
    tab_header_hex: str,
    palette: List[str],
):
    blank = (
        pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
    )
    slide = pres.slides.add_slide(blank)

    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_w = pres.slide_width - Inches(1.0)
    title_h = Inches(0.6)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = title_box.text_frame
    tf.text = f"{sheet_name} — {metric}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    slide_w, slide_h = pres.slide_width, pres.slide_height
    margin_l, margin_r, margin_bot = Inches(0.5), Inches(0.5), Inches(0.4)
    gap = Inches(0.25)

    usable_top = title_top + title_h + Inches(0.1)
    usable_left = margin_l
    usable_width = slide_w - (margin_l + margin_r)
    usable_height = slide_h - margin_bot - usable_top
    if usable_height <= Inches(1.0):
        return

    MIN_CHART_H = Inches(2.2)

    table_target_h = int(usable_height * 0.45)
    chart_target_h = usable_height - table_target_h - gap

    if chart_target_h < MIN_CHART_H:
        table_target_h = max(Inches(1.4), usable_height - MIN_CHART_H - gap)
        chart_target_h = usable_height - table_target_h - gap

    q3 = _build_q_summary_metric(pivot_q, metric)

    table_left = usable_left
    table_top = usable_top
    table_w = usable_width
    table_max_h = table_target_h

    used_table_h = _add_table_in_rect(
        slide=slide,
        df=q3,
        left=table_left,
        top=table_top,
        width=table_w,
        max_height=table_max_h,
        header_hex=tab_header_hex,
        min_body_row_h=Inches(0.22),
        header_row_h=Inches(0.28),
    )

    space_left_for_chart = usable_height - used_table_h - gap
    if space_left_for_chart < MIN_CHART_H and not q3.empty:
        header_h = Inches(0.28)
        min_body_h = Inches(0.22)
        allow_for_table = max(Inches(0.8), usable_height - MIN_CHART_H - gap)
        max_body_rows = max(1, int((allow_for_table - header_h) // int(min_body_h)))

        used_table_h = _add_table_in_rect(
            slide=slide,
            df=q3,
            left=table_left,
            top=table_top,
            width=table_w,
            max_height=allow_for_table,
            header_hex=tab_header_hex,
            min_body_row_h=min_body_h,
            header_row_h=header_h,
            prefer_rows=max_body_rows,
        )
        space_left_for_chart = usable_height - used_table_h - gap

    chart_left = usable_left
    chart_top = table_top + used_table_h + gap
    chart_w = usable_width
    chart_h = max(MIN_CHART_H, space_left_for_chart)

    _add_metric_chart(
        slide=slide,
        pivot_week=pivot_week,
        metric=metric,
        area_if_many=True,
        palette=palette,
        rect=(chart_left, chart_top, chart_w, chart_h),
    )


def _add_metric_vs_cost_slide(
    pres,
    sheet_name: str,
    pivot_week: pd.DataFrame,
    metric: str,
    cost_var: str,
    palette: List[str] | None,
    template_chart_name: str | None = None,
):
    """Add simple single-axis chart (metric vs cost) — NO dual axis."""
    if not cost_var or _norm(metric) == _norm(cost_var) or pivot_week.empty:
        return

    blank = (
        pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
    )
    slide = pres.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.text = f"{sheet_name} — {metric} vs {cost_var}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    date_col = pivot_week.columns[0]
    dates = pd.to_datetime(pivot_week[date_col], errors="coerce")
    cats = dates.dt.strftime("%b '%y").fillna("").tolist()

    def sum_or_total_list(df: pd.DataFrame, base: str) -> List[float]:
        tot_cols = [c for c in df.columns if _is_total_metric(c, base)]
        if tot_cols:
            return (
                pd.to_numeric(df[tot_cols[0]], errors="coerce")
                .fillna(0)
                .astype(float)
                .tolist()
            )

        cols = [c for c in df.columns[1:] if _ends_with_metric(c, base)]
        if not cols:
            return [0.0] * len(df)
        if len(cols) == 1:
            return (
                pd.to_numeric(df[cols[0]], errors="coerce")
                .fillna(0)
                .astype(float)
                .tolist()
            )

        s = (
            df[cols]
            .apply(pd.to_numeric, errors="coerce")
            .fillna(0)
            .sum(axis=1)
            .astype(float)
        )
        return s.tolist()

    y_metric = sum_or_total_list(pivot_week, metric)
    y_cost = sum_or_total_list(pivot_week, cost_var)
    cats, (y_metric, y_cost) = _align_cats_and_series(cats, [y_metric, y_cost])
    if not cats or (not any(y_metric) and not any(y_cost)):
        return

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series(metric, y_metric)
    chart_data.add_series(cost_var, y_cost)

    x, y, cx, cy = Inches(0.5), Inches(1.0), Inches(9.0), Inches(6.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    # --- Styling ---
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.has_major_gridlines = False

    pal = palette or []
    for idx, s in enumerate(chart.series):
        try:
            col = pal[idx % len(pal)] if pal else None
            if col and _is_hex_color(col):
                rgb = _hex_to_rgb(col)
                s.format.line.width = Pt(2)
                s.format.line.color.rgb = rgb
        except Exception:
            pass


def run(
    params: Dict, template_path: Optional[str], df: pd.DataFrame, output_path: str
) -> str:
    pres = Presentation(template_path) if template_path else Presentation()
    blank = (
        pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
    )

    if not template_path:
        pres.slides.add_slide(blank)
        pres.slides.add_slide(blank)

    sections = _compute_pivots(df, params)
    cost_var = (params.get("cost_var", "") or "").strip()
    tab_header_hex = (params.get("tab_color", "") or "").strip()
    dual_axis_chart_name = (
        params.get("dual_axis_chart_name", "") or "DualAxisChart"
    ).strip()

    for sec in sections:
        s_title = pres.slides.add_slide(blank)
        tb = s_title.shapes.add_textbox(
            Inches(0.5), pres.slide_height / 2 - Inches(0.5), Inches(9.0), Inches(1.0)
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"{sec['sheet_name']}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)
        p.font.bold = True

        for m in sec["metrics"]:
            _add_metric_slide(
                pres,
                sec["sheet_name"],
                sec["pivot_week"],
                sec["pivot_q"],
                metric=m,
                tab_header_hex=tab_header_hex,
                palette=sec["palette"],
            )

        if cost_var:
            for m in sec["metrics"]:
                _add_metric_vs_cost_slide(
                    pres,
                    sec["sheet_name"],
                    sec["pivot_week"],
                    metric=m,
                    cost_var=cost_var,
                    palette=sec["palette"],
                    template_chart_name=dual_axis_chart_name,
                )

    pres.save(output_path)
    return output_path
